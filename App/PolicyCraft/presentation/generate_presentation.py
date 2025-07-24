from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PolicyCraft"
    subtitle.text = "AI-Powered Policy Analysis for Higher Education\nJacek Robert Kszczot\nMSc Computer Science\nLeeds Trinity University"

    # Table of Contents
    toc_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(toc_slide_layout)
    title = slide.shapes.title
    title.text = "Table of Contents"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "1. Project Overview"
    p = tf.add_paragraph()
    p.text = "2. Key Features"
    p = tf.add_paragraph()
    p.text = "3. Technical Implementation"
    p = tf.add_paragraph()
    p.text = "4. Documentation & Compliance"
    p = tf.add_paragraph()
    p.text = "5. Demo"
    p = tf.add_paragraph()
    p.text = "6. Future Work"

    # Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "The Challenge:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Universities struggle with developing comprehensive AI policies"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Need for consistent, research-based recommendations"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Our Solution:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AI-powered policy analysis platform"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Evidence-based recommendations"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "• Academic citation validation"
    p.level = 2

    # Save the presentation
    prs.save('PolicyCraft_Presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
